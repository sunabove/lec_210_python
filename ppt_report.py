# coding=utf-8
import logging

logging.basicConfig( format='%(levelname)-8s %(asctime)s %(filename)s %(lineno)d %(message)s', level=logging.DEBUG )

class Report :
    def __init__(self) :
        pass
    pass

    def makeReport(self) :

        line = "\n" + "#"*80 + "\n"

        # import libraries
        from bs4 import BeautifulSoup 

        # specify the url
        html_url = "https://ko.wikipedia.org/wiki/%EA%B3%A0%EB%A0%A4%EC%9D%B8"

        # query the website and return the html to the variable ‘page’
        logging.info( "Getting a html data from %s" % html_url )

        from urllib.request import urlopen
        html_page = urlopen( html_url )
        html_src = html_page.read()

        logging.info( "Done. Getting a html data from %s" % html_url )

        # parse the html using beautiful soup and store in variable `soup`
        html = BeautifulSoup( html_src, "html.parser" , from_encoding="euc-kr" )

        # PPT 생성 
        from pptx import Presentation 
        from pptx.util import Pt
        prs = Presentation()

        # 첫번째 제목 PPT 슬라이드 추가 
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        logging.info( "title.string = %s" % html.h1.string ) 

        ppt_file_name =  html.h1.string

        title.text = html.h1.string
        subtitle.text = "2018-11-03\n한세대학교"

        # 두번째 개요 PPT 슬라이드 추가 
        bullet_slide_layout = prs.slide_layouts[1]

        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = "개요"

        tf = body_shape.text_frame 
        tf.text = html.h1.string
        tf.paragraphs[0].font.size = Pt(18)

        div = html.find( "div", { "class": "mw-parser-output" } )

        pList = div.find_all( "p" ) 

        for i in range( 2, 5 ) :
            p = tf.add_paragraph()
            p.text = pList[ i ].getText().strip()
            p.level = 1 
            p.font.size = Pt(14)
        pass

        # 세번째 목차 PPT 슬라이드 추가 
        bullet_slide_layout = prs.slide_layouts[1]

        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = "목차"

        tf = body_shape.text_frame
        div = html.find( "div", {"id": "toc" , "class": "toc" } )
        liList = div.find( "ul" ).find_all( "li" ) 

        for li in liList :
            p = tf.add_paragraph()
            p.text = li.getText().strip()
            p.level = 1 
            p.font.size = Pt(10)
        pass 

        # 네번째 역사 PPT 슬라이드 추가 
        bullet_slide_layout = prs.slide_layouts[1]

        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        div = html.find( "div", { "class": "mw-parser-output" } )

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = "역사"

        tf = body_shape.text_frame 

        # PPT 저장 
        prs.save( "%s.pptx" % ppt_file_name ) 

    pass

pass

if __name__ == '__main__' :   
    report = Report()
    report.makeReport()
pass
