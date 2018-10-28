# coding=utf-8

print( "Hello..." )

from pptx import Presentation
from pptx.util import Inches

# 이미지 가져와서 저장하기
import shutil
import requests

my_url = 'https://www.washingtonian.com/wp-content/uploads/2017/06/6-30-17-goat-yoga-congressional-cemetery-1-994x559.jpg'
response = requests.get(my_url, stream=True)
with open( 'my_image.png', 'wb') as file:
    shutil.copyfileobj(response.raw, file) 
pass

img_path = 'my_image.png'

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

left = top = Inches(1) 
height = Inches(5.5)
pic = slide.shapes.add_picture(img_path, left, top, height=height)

prs.save('test.pptx')

print( "Good bye!")